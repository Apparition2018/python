"""
packages：https://pypi.org/
"""

import json
import logging
import os
import random
import re
import tempfile
import time

import pandas
import requests
from fake_useragent import UserAgent
from tenacity import retry, stop_after_attempt

from util import Paths


class TestRequests:
    """
    `requests <https://pypi.org/project/requests/>`_：HTTP 库
    """

    def test_basic(self):
        r = requests.get('https://httpbin.org/basic-auth/user/pass', auth=('user', 'pass'))
        r.raise_for_status()
        assert r.status_code == 200
        assert r.headers['content-type'] == 'application/json'
        assert r.encoding == 'utf-8'
        # 响应内容（字节）
        assert r.content == b'{\n  "authenticated": true, \n  "user": "user"\n}\n'
        assert r.text == r.content.decode(r.encoding) == '{\n  "authenticated": true, \n  "user": "user"\n}\n'
        assert r.json() == json.loads(r.text) == {'authenticated': True, 'user': 'user'}

    def test_advance(self):
        from requests.adapters import HTTPAdapter
        from urllib3 import Retry
        # 会话
        with requests.Session() as session:
            # 连接池
            adapter = HTTPAdapter(
                pool_connections=15,
                pool_maxsize=50,
                # 重试
                max_retries=Retry(total=3, backoff_factor=0.5)
            )
            session.mount('https://', adapter)
            r = session.get(
                'https://2025.ip138.com/',
                headers={'User-Agent': UserAgent().random},
                # 代理：https://free.kuaidaili.com/free/dps/
                # proxies={"https": "180.121.147.240:20756"},
                timeout=2,
                allow_redirects=True
            )
            print(re.search(r'<title[^>]*>.*?(\d+(?:\.\d+)*)', r.text).group(1).strip())


def test_user_agent():
    """
    `fake-useragent <https://pypi.org/project/fake-useragent/>`_：useragent faker
    """
    from fake_useragent import UserAgent
    url = 'https://www.baidu.com'
    r = requests.get(url)
    # 1. 自定义随机 User-Agent
    user_agents = [
        'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/127.0.0.0 Safari/537.36',
        'Mozilla/5.0 (iPhone; CPU iPhone OS 16_6 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.6 Mobile/15E148 Safari/604.1',
        'Mozilla/5.0 (Linux; Android 10; K) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/127.0.0.0 Mobile Safari/537.36'
    ]
    r2 = requests.get(url, headers={'User-Agent': random.choice(user_agents)})
    assert len(r.content) < len(r2.content)
    # 2. 使用 fake_useragent 随机 User-Agent
    r3 = requests.get(url, headers={'User-Agent': UserAgent().random})
    assert len(r.content) < len(r3.content)


class TestSelenium:
    """
    `selenium <https://pypi.org/project/selenium/>`_：自动化 web 浏览器交互
    """

    def test_selenium(self):
        from selenium import webdriver
        from selenium.webdriver.chrome.service import Service
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as ec
        options = webdriver.ChromeOptions()
        # 不显示图形界面
        # options.add_argument('headless')
        service = Service(executable_path=str(Paths.fixture('chromedriver.exe')))
        driver = webdriver.Chrome(options=options, service=service)
        driver.get('https://www.baidu.com')

        print(driver.current_url)
        print(driver.title)
        # print(webdriver.page_source)
        # 截图
        with tempfile.NamedTemporaryFile() as fp:
            driver.save_screenshot(fp.name)
        # region 搜索 selenium
        chat_textarea = driver.find_element(By.ID, 'chat-textarea')
        chat_textarea.send_keys('selenium')
        chat_submit_button = driver.find_element(By.ID, 'chat-submit-button')
        chat_submit_button.click()
        # 等待直到，某元素可见
        WebDriverWait(driver, 10).until(ec.visibility_of_element_located((By.ID, 'page')))
        # endregion
        # 通过执行 js 滚动到底部
        driver.execute_script('window.scrollTo(0, document.documentElement.scrollHeight);')
        time.sleep(1)
        # 通过执行 js 打开一个窗口
        driver.execute_script('window.open("https://www.sogou.com")')
        # 所有窗口句柄
        handles = driver.window_handles
        # 切换窗口
        driver.switch_to.window(handles[1])
        time.sleep(1)
        # 关闭浏览器并关闭 ChromeDriver
        driver.quit()


class TestTenacity:
    """
    `tenacity <https://pypi.org/project/tenacity/>`_：多用途重试库
    """

    @staticmethod
    def do_something():
        if random.random() < 0.8:
            print('失败')
            raise Exception('异常')
        else:
            print('成功')

    @staticmethod
    def do_something2():
        if random.random() < 0.7:
            print('失败')
            return False
        else:
            print('成功')
            return True

    def test_stop(self):
        # stop_after_attempt：最多重试次数
        from tenacity import stop_after_delay, stop_never, stop_any, stop_all
        do_something = retry(stop=stop_after_attempt(10))(self.do_something)
        do_something()
        # stop_after_delay：最多重试总秒数
        do_something = retry(stop=stop_after_delay(10))(self.do_something)
        do_something()
        # stop_never：重试直到成功；不设置 stop，效果等同于 stop_never
        do_something = retry(stop=stop_never)(self.do_something)
        do_something()
        # stop_any
        do_something = retry(stop=stop_any(stop_after_attempt(5), stop_after_delay(5)))(self.do_something)
        do_something()
        # stop_all
        do_something = retry(stop=stop_all(stop_after_attempt(5), stop_after_delay(5)))(self.do_something)
        do_something()

    def test_wait(self):
        from tenacity import wait_fixed, wait_random, wait_exponential, wait_random_exponential, wait_none
        # wait_fixed：每次等待秒数
        do_something = retry(wait=wait_fixed(1))(self.do_something)
        do_something()
        # wait_random：每次等待随机秒数
        do_something = retry(wait=wait_random(1, 3))(self.do_something)
        do_something()
        # wait_exponential：指数退避，等待 multiplier*2^尝试次数，即依次等待 1 2 4 8 10 10 10 …
        do_something = retry(wait=wait_exponential(1, 10))(self.do_something)
        do_something()
        # wait_random_exponential：等待 random()*multiplier*2^尝试次数
        do_something = retry(wait=wait_random_exponential(2, 10))(self.do_something)
        do_something()
        # wait_none：不等待
        do_something = retry(wait=wait_none())(self.do_something)
        do_something()

    def test_retry(self):
        from tenacity import (retry_if_exception, retry_if_exception_type, retry_if_result, retry_if_exception_message,
                              retry_any, retry_all)
        # retry_if_exception_type：只对指定异常重试
        do_something = retry(retry=retry_if_exception_type(Exception))(self.do_something)
        do_something()
        # retry_if_exception：抛出异常时执行 predicate，返回 True 重试
        do_something = retry(retry=retry_if_exception(lambda e: type(e) is Exception))(self.do_something)
        do_something()
        # retry_if_result：有返回值时执行 predicate，返回 True 重试
        do_something = retry(retry=retry_if_result(lambda res: res is False))(self.do_something2)
        do_something()
        # retry_if_exception_message：根据异常消息是否包含指定字符串决定是否重试
        do_something = retry(retry=retry_if_exception_message(match='异常'))(self.do_something)
        do_something()
        # retry_any
        do_something = retry(
            retry=retry_any(retry_if_exception_type(Exception), retry_if_exception_message(match='异常'))
        )(self.do_something)
        do_something()
        # retry_all
        do_something = retry(
            retry=retry_all(retry_if_exception_type(Exception), retry_if_exception_message(match='异常'))
        )(self.do_something)
        do_something()

    def test_before_after(self):
        """
        before → 失败 → after → before_sleep → 等待 → …
        """
        from tenacity import before_sleep_log
        def log_before(retry_state):
            print(f'🟢 [before] 第{retry_state.attempt_number}次尝试')

        def log_after(retry_state):
            print(f'🔴 [after]  总耗时{retry_state.seconds_since_start}秒\n')

        do_something = retry(
            before=log_before, after=log_after, before_sleep=before_sleep_log(logging.getLogger(), logging.INFO) # type: ignore
        )(self.do_something)
        do_something()

    def test_raise(self):
        from tenacity import RetryError
        # reraise=False，抛出 RetryError；reraise=True，抛出原始异常
        do_something = retry(stop=stop_after_attempt(1), reraise=True)(self.do_something)
        try:
            do_something()
        except RetryError:
            print('RetryError')
        except Exception:
            print('Exception')


def test_jsonpath_ng():
    """
    `jsonpath-ng <https://pypi.org/project/jsonpath-ng/>`_
    """
    from jsonpath_ng import parse
    r = requests.get('https://reqres.in/api/users', headers={'x-api-key': 'reqres-free-v1'})
    titles = [match.value for match in parse('$.data[*].first_name').find(r.json())]
    print(titles[:5])


def test_lxml():
    """
    `lxml <https://pypi.org/project/lxml/>`_：

    1. 通过封装 C 库（libxml2/libxslt）提供强大的 XML/HTML 处理能力
    2. 以 ElementTree API（xml.etree.ElementTree） 形式暴露给开发者，兼顾安全性与易用性
    3. 支持 XPath、RelaxNG、XML Schema、XSLT、C14N 等功能
    """
    from lxml import etree
    r = requests.get('https://www.w3schools.com/xml/books.xml')
    root = etree.XML(r.content)
    assert (root.xpath('/bookstore/book[2]/author/text()'), 'J K. Rowling')
    assert (root.xpath('/bookstore/book[2]/title/@lang/text()'), 'Harry Potter')


def test_moviepy():
    """
    `moviepy <https://pypi.org/project/moviepy/>`_：用于视频编辑：剪切、连接、插入标题、视频合成（也称为非线性编辑）、视频处理和创建自定义效果
    """
    import tempfile
    from moviepy import AudioFileClip, VideoFileClip
    ua = UserAgent()

    def get_windows_ua():
        return next((ua.random for _ in range(10) if 'Windows' in ua.random),
                    'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36')

    url = 'https://www.bilibili.com/video/BV1D4411L7Qd/'
    headers = {'User-Agent': get_windows_ua(), 'Referer': url}
    r = requests.get(url, headers=headers)
    audio_url = re.search(r'"id":30216,"baseUrl":"(.*?)","base_url"', r.text).group(1)
    video_url = re.search(r'"id":16,"baseUrl":"(.*?)","base_url"', r.text).group(1)
    audio_data = requests.get(audio_url, headers=headers).content
    video_data = requests.get(video_url, headers=headers).content
    with (tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as temp_audio,
          tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as temp_video):
        temp_audio.write(audio_data)
        temp_video.write(video_data)
        temp_audio_path, temp_video_path = temp_audio.name, temp_video.name
    try:
        with AudioFileClip(temp_audio_path) as audio, VideoFileClip(temp_video_path) as video:
            video.with_audio(audio).write_videofile(
                'video.mp4',
                threads=os.cpu_count(),
                codec='libx264',
                audio_codec='aac',
                preset='fast',
                ffmpeg_params=['-movflags', '+faststart']
            )
    finally:
        os.remove(temp_audio_path) if os.path.exists(temp_audio_path) else None
        os.remove(temp_video_path) if os.path.exists(temp_video_path) else None


def test_pymysql():
    """
    `PyMySQL <https://pypi.org/project/PyMySQL/>`_
    """
    import pymysql
    from pymysql import cursors
    connection = pymysql.connect(host='43.136.102.115',
                                 user='root',
                                 password='Cesc123!',
                                 database='epitome',
                                 charset='utf8mb4',
                                 cursorclass=cursors.DictCursor)
    with connection:
        with connection.cursor() as cursor:
            sql = 'INSERT INTO `demo` (`id`, `name`) VALUES (%s, %s)'
            cursor.execute(sql, ('1', 'ljh'))
        connection.commit()
        with connection.cursor() as cursor:
            sql = 'SELECT `id`, `name` FROM `demo` WHERE `id`=%s'
            cursor.execute(sql, 1)
            result = cursor.fetchone()
            print(result)
        with connection.cursor() as cursor:
            sql = 'DELETE FROM `demo` WHERE id=%s'
            cursor.execute(sql, 1)
        connection.commit()


class TestOpenpyxl:
    """
    `openpyxl <https://pypi.org/project/openpyxl/>`_：用于读取/写入Excel 2010 xlsx/xlsm/xltx/xltm 文件
    """
    import openpyxl
    XLSX_PATH = Paths.fixture('test.xlsx')

    def test_read(self):
        workbook = self.openpyxl.load_workbook(self.XLSX_PATH)
        # 工作表列表
        assert workbook.sheetnames == [sheet.title for sheet in workbook]
        # 工作表
        assert workbook.active == workbook[workbook.sheetnames[0]]
        worksheet = workbook.active
        # (行数，列数)
        assert (worksheet.max_row, worksheet.max_column) == (9, 3)
        # 工作表切片
        assert worksheet['A1:B2'] == ((worksheet['A1'], worksheet['B1']), (worksheet['A2'], worksheet['B2']))
        # 行切片
        assert worksheet['2:3'] == (worksheet[2], worksheet[3])
        # 列切片
        assert worksheet['C:D'] == (worksheet['C'], worksheet['D'])
        # 单元格
        assert worksheet['A'][0] == worksheet['A1'] == worksheet.cell(1, 1)
        # 遍历工作表
        for worksheet in workbook:
            # 遍历行
            for row in worksheet:
                assert isinstance(row, tuple)
                for cell in row:
                    print(cell.coordinate, cell.row, cell.column, cell.value)
            # 遍历行
            for row in range(1, worksheet.max_row + 1): pass
            # 遍历行
            for row in worksheet.iter_rows(min_row=1, max_row=2, min_col=1, max_col=2): pass
            # 遍历列
            for column in worksheet.columns: pass
        # 十进制 ↔ 列字母
        from openpyxl.utils import get_column_letter, column_index_from_string
        assert (get_column_letter(100), column_index_from_string('CV')) == ('CV', 100)

    def test_write(self):
        from openpyxl.utils import get_column_letter
        from openpyxl.styles import Font
        workbook = self.openpyxl.load_workbook(self.XLSX_PATH)
        worksheet = workbook[workbook.sheetnames[1]]
        # 设置工作表标题
        worksheet.title = 'write'
        for _ in range(worksheet.max_row + 1):
            # 删除行
            worksheet.delete_rows(1)
        for row in range(1, 10):
            # 在工作表底部追加一行值
            worksheet.append([f'{get_column_letter(col)}{row}' for col in range(1, 10)])
            # 设置行高
            worksheet.row_dimensions[row].height = 25
            for col in range(1, 10):
                # 设置列宽
                if row == 1: worksheet.column_dimensions[get_column_letter(col)].width = 4.78
                # 设置单元格字体
                worksheet[f'{get_column_letter(col)}{row}'].font = Font(size=14, italic=True)
        row = 10
        for col in range(1, 10):
            # 设置单元格
            worksheet.cell(row, col, f'{get_column_letter(col)}{row}')
        # 合并单元格
        worksheet.merge_cells('A11:B12')
        worksheet['A11'] = 'A11'
        # 删除工作表
        workbook.remove(workbook[workbook.sheetnames[2]])
        # 创建工作表
        workbook.create_sheet('new', 2)
        # 保存当前工作簿
        workbook.save(self.XLSX_PATH)


class TestPythonDocx:
    """
    `python-docx <https://pypi.org/project/python-docx/>`_：用于读取、创建和更新 Microsoft Word 2007+（.docx）文件
    """
    from docx import Document
    DOCX_PATH = str(Paths.fixture('test.docx'))

    def test_read(self):
        from docx.enum.style import WD_STYLE_TYPE
        doc = type(self).Document(self.DOCX_PATH)
        # 遍历段落
        for paragraph in doc.paragraphs:
            # 所有标题
            if re.match(r'^Heading \d+$', paragraph.style.name): pass
        # 遍历样式
        for style in doc.styles:
            if style.type == WD_STYLE_TYPE.PARAGRAPH:
                # 所有标题
                if re.match(r'^Heading \d+$', style.name): pass
                # 所有正文
                if style.name == 'Normal': pass

    def test_write(self):
        from io import BytesIO
        from docx.enum.section import WD_ORIENTATION
        from docx.enum.text import WD_BREAK
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        from docx.shared import Inches
        from docx.shared import Pt
        doc = type(self).Document(self.DOCX_PATH)
        section = doc.sections[0]
        # 页面宽高
        assert round(section.page_width.cm, 1) == 21.0
        assert round(section.page_height.cm, 1) == 29.7
        # 页面方向
        assert section.orientation == WD_ORIENTATION.PORTRAIT
        body = doc.element.body
        # 清空文档（段落、表格、浮动文本框；图片、分页符都是在段落内）
        for e in list(body):
            if not e.tag.endswith('}sectPr'):
                body.remove(e)
        # region 添加内容
        # 添加标题（应用了 Heading \d 样式的段落）
        doc.add_heading('一级标题', 1)
        doc.add_heading('二级标题', 2)
        doc.add_heading('三级标题', 3)
        # 添加段落
        p = doc.add_paragraph('正文')
        # 间距
        p.paragraph_format.space_before = Pt(6)
        p.paragraph_format.space_after = Pt(6)
        # 首行缩进
        p.paragraph_format.first_line_indent = Pt(21)
        p.add_run('普通')
        p.add_run('加粗').bold = True
        p.add_run('斜体').italic = True
        # 插入段落
        doc.paragraphs[3].insert_paragraph_before('插入')
        doc.paragraphs[3].alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # 添加图片（在段落内）
        img_bio = BytesIO(requests.get('https://i.imgs.ovh/2025/04/18/j9kHY.jpeg').content)
        doc.add_picture(img_bio, width=Inches(1.25))
        # 相当于
        doc.add_paragraph().add_run().add_picture(img_bio, width=Inches(1.25))
        # 水平居中
        # 添加表格
        rows = cols = 3
        doc.add_table(rows, cols)
        for i in range(rows):
            for j in range(cols):
                doc.tables[0].cell(i, j).paragraphs[0].add_run(str(i * cols + j + 1))
        # 添加分页符（在段落内）
        doc.add_page_break()
        # 相当于
        doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        # endregion
        doc.save(self.DOCX_PATH)


def test_pandas():
    """
    `pandas <https://pypi.org/project/pandas/>`_：数据分析工具包

    处理数据：避免遍历操作

    1. 数学计算：dataframe['a'] * 2
    2. 条件判断：numpy.where() 或 dataframe.loc[dataframe['a'] > 30, 'b']
    3. 多条件：numpy.select()
    4. 分组：groupby()
    5. 复杂多级：apply()
    """
    excel_path = Paths.fixture('test.xlsx')
    # 读取工作表，以下三行均表示读取第一个工作表
    dataframe = pandas.read_excel(excel_path)
    dataframe = pandas.read_excel(excel_path, sheet_name='read')
    dataframe = pandas.read_excel(excel_path, sheet_name=0)
    # 读取多个工作表
    dataframes = pandas.read_excel(excel_path, sheet_name=['read', 2])
    # 读取特定列
    dataframe = pandas.read_excel(excel_path, usecols=['caseid', 'data'])
    dataframe = pandas.read_excel(excel_path, usecols=[0, 2])
    # 其它
    dataframe = pandas.read_excel(excel_path,
                                  # 跳过n行再读取n行
                                  skiprows=0, nrows=9,
                                  # 指定某列数据类型
                                  dtype={'caseid': int, 'excepted': str, 'data': str},
                                  # 将某些值识别为 NaN
                                  na_values=['无', '/', ''],
                                  # 自动解释日期列
                                  # parse_dates=['date'],
                                  # 指定日期格式
                                  # date_format='%Y-%m-%d'
                                  )
    print(dataframe.info())
    # 删除重复数据
    dataframe.drop_duplicates()
    # 遍历数据
    for row in dataframe.itertuples():
        print(f'Index: {row.Index}, caseid: {row.caseid}, excepted: {row.excepted}, data: {row.data}')
    # 两种索引方式读取数据
    assert dataframe.iloc[0, 0] == dataframe.loc[0, 'caseid'] == 1
    with tempfile.TemporaryDirectory() as tmpdir:
        # 写入 csv 文件
        dataframe.to_csv(f'{tmpdir}/test.csv', encoding='utf-8')
        # 写入 Excel 工作表
        dataframe.to_excel(f'{tmpdir}/test2.xlsx')


def test_pdfplumber():
    """
    `pdfplumber <https://pypi.org/project/pdfplumber/>`_：获取 PDF 每个 char、rectangle、line 的信息
    """
    import pdfplumber
    with pdfplumber.open(Paths.fixture('test.pdf')) as pdf:
        for page in pdf.pages:
            print(page.extract_text)
            for table in page.extract_tables():
                print(table)


def test_pypdf():
    """
    `pypdf <https://pypi.org/project/pypdf/>`_：拆分、合并、裁剪和转换 PDF 文件
    """
    from pypdf import PdfReader, PdfWriter
    pwd = '123456'
    reader = PdfReader(Paths.fixture('test.pdf'))
    # 取消密码
    reader.decrypt(pwd)
    # watermark_page = PdfReader('watermark.pdf').get_page(0)
    with tempfile.TemporaryFile() as fp:
        writer = PdfWriter()
        # 设置密码
        writer.encrypt(pwd)
        for i, page in enumerate(reader.pages):
            # 给每一页合并水印
            # page.merge_page(watermark_page)
            writer.add_page(page)
        writer.write(fp)


class TestPPTX:
    """
    `python-pptx <https://pypi.org/project/python-pptx/>`_：创建、读取和更新 PowerPoint (.pptx) 文件
    """
    from pptx import Presentation
    PPTX_PATH = str(Paths.fixture('test.pptx'))

    def test_read(self):
        # 读取演示文稿
        p = type(self).Presentation(self.PPTX_PATH)
        # 遍历幻灯片
        for slide in p.slides:
            # 遍历 shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        print(paragraph.text)
                if shape.is_placeholder:
                    print(f'名称：{shape.name}，索引：{shape.placeholder_format.idx}，类型：{shape.placeholder_format.type}')

    def test_write(self):
        from pptx.chart.data import ChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
        from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
        from pptx.util import Cm, Pt
        # 新建演示文稿
        p = type(self).Presentation()
        # 添加幻灯片，使用布局索引为 0 的版式
        # slide_layouts 是一个包含 11 种默认版式的列表
        slide = p.slides.add_slide(p.slide_layouts[0])
        # 遍历 shapes，找到类型为 TITLE 的第一个 is_placeholder 设置文本
        slide.shapes.title.text_frame.text = 'Title'
        # 设置第二个 placeholder 的文本
        slide.placeholders[1].text = 'Placeholder'
        # 获取第二个 placeholder 的 text_frame
        text_frame = slide.placeholders[1].text_frame
        # 添加一个段落到 text_frame，并设置其层级为 1
        paragraph = text_frame.add_paragraph()
        paragraph.text = 'Paragraph'
        paragraph.level = 1
        # 添加一个新段落到 text_frame，并设置其层级为 2
        paragraph2 = text_frame.add_paragraph()
        paragraph2.text = 'Paragraph2'
        paragraph2.level = 2
        # 添加一个 run 到 paragraph
        run = paragraph2.add_run()
        run.text = 'Run'
        run.font.size = Pt(18)
        # 添加一个 textbox
        textbox = slide.shapes.add_textbox(left=Cm(1), top=Cm(1), width=Cm(2.5), height=Cm(1))
        textbox.text_frame.text = 'Textbox'
        textbox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(100, 100, 100)
        textbox.line.color.rgb = RGBColor(0, 0, 0)
        textbox.line.width = Cm(0.1)
        # 添加一个 table
        table = slide.shapes.add_table(rows=3, cols=4, left=Cm(5), top=Cm(1), width=Cm(6), height=Cm(4.5)).table
        # 添加一个 chart
        chart_data = ChartData()
        chart_data.categories = ['1', '2', '3', '4']
        chart_data.add_series(name='销量', values=[5676, 4563, 7656, 8685])
        chart_data.add_series(name='金额', values=[3246, 2436, 6343, 3565])
        chart = slide.shapes.add_chart(chart_type=XL_CHART_TYPE.BAR_CLUSTERED, x=Cm(12), y=Cm(1), cx=Cm(12), cy=Cm(6),
                                       chart_data=chart_data).chart
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.number_format = '#,#'
        chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END
        p.save(self.PPTX_PATH)
